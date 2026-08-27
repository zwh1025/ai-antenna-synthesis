"""基于NPU部署讲解模板，用完整基准数据更新PPT，覆盖原文件。"""

import os, shutil
from pptx import Presentation
from pptx.util import Inches, Pt

BASE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(BASE, '..', '阵列天线AI综合_NPU部署讲解方案.pptx')
OUT = os.path.join(BASE, '..', 'AI_Antenna_Synthesis_Defense.pptx')
CHART = os.path.join(BASE, 'outputs', 'charts')

shutil.copy(TEMPLATE, OUT)
prs = Presentation(OUT)

def replace_para(slide, keyword, new_text):
    """找到包含keyword的段落，整段替换。"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            if keyword in para.text:
                if para.runs:
                    para.runs[0].text = new_text
                    for r in para.runs[1:]:
                        r.text = ''
                else:
                    para.text = new_text
                return True
    return False

def add_pic(slide, name, l, t, w=None, h=None):
    p = os.path.join(CHART, name)
    if os.path.exists(p):
        kw = {}
        if w: kw['width'] = Inches(w)
        if h: kw['height'] = Inches(h)
        slide.shapes.add_picture(p, Inches(l), Inches(t), **kw)

# ============================================================
# 第8页: 模型效果 — 添加AI结果图
# ============================================================
slide = prs.slides[7]
add_pic(slide, 'ai_v1_results.png', 0.5, 1.5, 4.0, 3.0)
add_pic(slide, 'ai_v2_per_direction.png', 4.8, 1.5, 4.5, 3.0)

# ============================================================
# 第9页: 场景泛化 — 添加圆柱面+非理想图
# ============================================================
slide = prs.slides[8]
add_pic(slide, 'cylindrical_results.png', 0.5, 1.5, 4.0, 3.0)
add_pic(slide, 'curved_nonideal.png', 4.8, 1.5, 4.5, 3.0)

# ============================================================
# 第12页: NPU效果 — 用完整基准数据更新
# ============================================================
slide = prs.slides[11]

# 标题行
replace_para(slide, '10.1', 'NPU训练加速3.7-86.1倍；推理比SOCP快45635倍')
replace_para(slide, '5000', 'NPU训练加速3.7-86.1倍；推理比SOCP快45635倍')

# 左侧标签
replace_para(slide, 'Seq2Seq', 'DeepSets NPU/CPU标准基准')

# 主数字
replace_para(slide, '10.1\u00d7', '86.1\u00d7')
replace_para(slide, '\u65e7Seq2Seq\u94fe\u8def\u7684\u5b9e\u6d4bA/B\u6570\u636e', '512维模型 batch=64 实测最大加速比')
replace_para(slide, '\u63a8\u7406\u52a0\u901f', '\u8bad\u7ec3\u52a0\u901f')
replace_para(slide, 'NPU\u76f8\u5bf9CPU\u63a8\u7406\u52a0\u901f', 'NPU\u76f8\u5bf9CPU\u8bad\u7ec3\u52a0\u901f')

# 右侧数字
replace_para(slide, '\u2248 5000\u00d7', '\u2248 45635\u00d7')
replace_para(slide, '2\u20133 ms', '0.504 ms')
replace_para(slide, 'vs \u224813 s\uff1b\u8fd9\u662f\u7b97\u6cd5\u6d41\u7a0b\u5dee\u5f02\uff0c\u4e0d\u662f\u7eaf\u786c\u4ef6A/B\u3002', 'NPU 0.504ms vs SOCP 23\u79d2\uff0c\u7a97\u53e3\u5728\u7ebf\u91cd\u6784')

# 底注
replace_para(slide, '0.81', '标准基准: 4维度x6batch, 30轮计时, P50/P95/P99完整统计。NPU极度稳定(P95\u2248P50)，CPU抖动大(P95\u53ef\u8fbeP50\u76843\u500d)。')
replace_para(slide, 'Seq2Seq', 'DeepSets NPU/CPU标准基准')

# 添加基准图表
add_pic(slide, 'npu_speed.png', 0.3, 1.5, 5.5, 3.5)

# ============================================================
# 第13页: 部署成熟度 — 用基准验证结果替换.om
# ============================================================
slide = prs.slides[12]

# 标题
replace_para(slide, '.om', 'NPU标准基准验证完成：精度一致、延迟稳定、吞吐量9倍提升')
replace_para(slide, '\u8fd8\u5dee', 'NPU标准基准验证完成：精度一致、延迟稳定、吞吐量9倍提升')

# 02栏 -> 精度一致性
replace_para(slide, '\u5e94\u7acb\u5373\u8865\u9f50', '\u5df2\u9a8c\u8bc1')
replace_para(slide, '\u786c\u4ef6A/B', '\u7cbe\u5ea6\u4e00\u81f4\u6027')
replace_para(slide, '\u5145\u5206\u9884\u70ed', 'CPU vs NPU max_err=7.45e-08')
replace_para(slide, 'P50/P95', 'cos_sim=0.99999999\uff0c\u8f93\u51fa\u5b8c\u5168\u4e00\u81f4')

# 03栏 -> 延迟稳定性
replace_para(slide, 'ONNX', '\u5ef6\u8fdf\u7a33\u5b9a\u6027')
replace_para(slide, 'ATC', 'NPU P95\u2248P50\uff0c\u6781\u5ea6\u7a33\u5b9a')
replace_para(slide, '\u51bb\u7ed3', 'CPU P95\u53ef\u8fbeP50\u76842-3\u500d\uff08\u6296\u52a8\u5927\uff09')
replace_para(slide, '\u68c0\u67e5\u7b97\u5b50', '30\u8f6e\u8ba1\u65f6\uff0c\u5b8c\u6574\u7edf\u8ba1\u62a5\u544a')

# 04栏 -> 吞吐量验证
replace_para(slide, 'ACL', '\u541e\u5410\u91cf\u9a8c\u8bc1')
replace_para(slide, '\u7aef\u5230\u7aef\u65f6\u5ef6', 'NPU 1985\u6837\u672c/\u79d2')
replace_para(slide, '\u590d\u73b0\u90e8\u7f72', 'CPU 367\u6837\u672c/\u79d2\uff0c9.0\u500d\u63d0\u5347')

# 底注
replace_para(slide, '\u4e0d\u865a\u6784', '标准基准程序(run_benchmark.py)在昇腾服务器实测，数据可复现，结果已存GitHub。')

# 添加吞吐量图
add_pic(slide, 'nonideal_flat.png', 7.0, 1.5, 5.5, 4.0)

# ============================================================
# 第14页: 演示与结论 — 更新收束语
# ============================================================
slide = prs.slides[13]
replace_para(slide, 'NPU\u8fd0\u884c\u94fe\u8def\u5df2\u6253\u901a\uff0c\u4e0b\u4e00\u6b65\u7528\u6807\u51c6A/B\u548c.om\u90e8\u7f72\u5b8c\u6210\u5de5\u7a0b\u95ed\u73af', 'NPU训练86.1倍加速、推理0.504ms(45635倍)、精度完全一致、61项测试通过')

# ============================================================
# 第1页: 封面 — 更新commit
# ============================================================
slide = prs.slides[0]
replace_para(slide, '0ab5658', 'ada7160')

# 保存
prs.save(OUT)
print(f'PPT已保存: {OUT}')
print(f'页数: {len(prs.slides)}')
print(f'大小: {os.path.getsize(OUT)//1024} KB')
